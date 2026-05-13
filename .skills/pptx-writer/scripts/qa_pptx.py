#!/usr/bin/env python3
"""Run conservative package/content QA checks on a PPTX file."""

from __future__ import annotations

import argparse
import json
import re
import zipfile
from dataclasses import asdict, dataclass
from pathlib import Path
from typing import Any

from pptx import Presentation


PLACEHOLDER_PATTERNS = [
    re.compile(pattern, re.IGNORECASE)
    for pattern in [
        r"click to add",
        r"slide number",
        r"\bsldNum\b",
        r"lorem ipsum",
        r"\bplaceholder\b",
        r"\bTODO\b",
    ]
]


@dataclass(frozen=True)
class Finding:
    severity: str
    slide: int | None
    message: str


def _extract_slide_text(slide: Any) -> str:
    texts: list[str] = []
    for shape in slide.shapes:
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    text = cell.text.strip()
                    if text:
                        texts.append(text)
        if getattr(shape, "has_text_frame", False):
            text = shape.text.strip()
            if text:
                texts.append(text)
    return "\n".join(texts)


def _package_text(path: Path) -> str:
    chunks: list[str] = []
    with zipfile.ZipFile(path) as zf:
        for name in zf.namelist():
            if name.startswith("ppt/") and name.endswith(".xml"):
                chunks.append(zf.read(name).decode("utf-8", errors="ignore"))
    return "\n".join(chunks)


def qa_pptx(path: Path) -> dict[str, Any]:
    findings: list[Finding] = []
    prs = Presentation(str(path))
    slide_summaries = []

    if not prs.slides:
        findings.append(Finding("error", None, "Deck has no slides."))

    for index, slide in enumerate(prs.slides, start=1):
        text = _extract_slide_text(slide)
        slide_summaries.append(
            {
                "slide": index,
                "layout": slide.slide_layout.name,
                "text_length": len(text),
                "text_preview": text[:300],
            }
        )
        if not text.strip():
            findings.append(Finding("warn", index, "Slide has no extractable text."))
        for pattern in PLACEHOLDER_PATTERNS:
            if pattern.search(text):
                findings.append(
                    Finding("error", index, f"Visible text matches placeholder pattern: {pattern.pattern}")
                )

    package_text = _package_text(path)
    for pattern in PLACEHOLDER_PATTERNS:
        if pattern.search(package_text):
            findings.append(
                Finding("warn", None, f"Package XML matches placeholder pattern: {pattern.pattern}")
            )

    errors = [finding for finding in findings if finding.severity == "error"]
    return {
        "path": str(path),
        "slide_count": len(prs.slides),
        "ok": not errors,
        "findings": [asdict(finding) for finding in findings],
        "slides": slide_summaries,
    }


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("pptx", type=Path)
    parser.add_argument("--output", type=Path)
    args = parser.parse_args()

    payload = qa_pptx(args.pptx)
    text = json.dumps(payload, ensure_ascii=False, indent=2)
    if args.output:
        args.output.parent.mkdir(parents=True, exist_ok=True)
        args.output.write_text(text + "\n", encoding="utf-8")
    else:
        print(text)
    return 0 if payload["ok"] else 1


if __name__ == "__main__":
    raise SystemExit(main())
