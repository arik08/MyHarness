#!/usr/bin/env python3
"""Extract a structural text inventory from a PPTX file."""

from __future__ import annotations

import argparse
import json
import zipfile
from pathlib import Path
from typing import Any

from pptx import Presentation


def _paragraphs(shape: Any) -> list[str]:
    if getattr(shape, "has_table", False):
        texts: list[str] = []
        for row in shape.table.rows:
            for cell in row.cells:
                text = cell.text.strip()
                if text:
                    texts.append(text)
        return texts
    if not getattr(shape, "has_text_frame", False):
        return []
    result: list[str] = []
    for paragraph in shape.text_frame.paragraphs:
        text = "".join(run.text for run in paragraph.runs).strip()
        if text:
            result.append(text)
    return result


def _slide_payload(slide: Any, index: int) -> dict[str, Any]:
    shapes = []
    for shape in slide.shapes:
        texts = _paragraphs(shape)
        shapes.append(
            {
                "shape_id": shape.shape_id,
                "name": shape.name,
                "shape_type": str(shape.shape_type),
                "is_placeholder": bool(getattr(shape, "is_placeholder", False)),
                "texts": texts,
            }
        )
    return {
        "slide_number": index + 1,
        "layout_name": slide.slide_layout.name,
        "text": "\n".join(text for shape in shapes for text in shape["texts"]),
        "shapes": shapes,
    }


def _package_summary(path: Path) -> dict[str, Any]:
    with zipfile.ZipFile(path) as zf:
        names = zf.namelist()
    return {
        "slide_xml_count": sum(
            1 for name in names if name.startswith("ppt/slides/slide") and name.endswith(".xml")
        ),
        "notes_xml_count": sum(1 for name in names if name.startswith("ppt/notesSlides/")),
        "media_count": sum(1 for name in names if name.startswith("ppt/media/")),
        "chart_count": sum(1 for name in names if name.startswith("ppt/charts/")),
    }


def read_pptx(path: Path) -> dict[str, Any]:
    prs = Presentation(str(path))
    return {
        "path": str(path),
        "slide_count": len(prs.slides),
        "package": _package_summary(path),
        "slides": [_slide_payload(slide, index) for index, slide in enumerate(prs.slides)],
    }


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("pptx", type=Path)
    parser.add_argument("--output", type=Path)
    args = parser.parse_args()

    payload = read_pptx(args.pptx)
    text = json.dumps(payload, ensure_ascii=False, indent=2)
    if args.output:
        args.output.parent.mkdir(parents=True, exist_ok=True)
        args.output.write_text(text + "\n", encoding="utf-8")
    else:
        print(text)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
