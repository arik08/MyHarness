#!/usr/bin/env python3
"""Inspect a PPTX template's masters, layouts, and placeholders."""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any

from pptx import Presentation


def _emu(value: int) -> float:
    return round(value / 914400, 4)


def _shape_payload(shape: Any) -> dict[str, Any]:
    placeholder = None
    if getattr(shape, "is_placeholder", False):
        fmt = shape.placeholder_format
        placeholder = {
            "idx": fmt.idx,
            "type": str(fmt.type),
        }
    return {
        "shape_id": shape.shape_id,
        "name": shape.name,
        "shape_type": str(shape.shape_type),
        "is_placeholder": bool(getattr(shape, "is_placeholder", False)),
        "placeholder": placeholder,
        "left_in": _emu(shape.left),
        "top_in": _emu(shape.top),
        "width_in": _emu(shape.width),
        "height_in": _emu(shape.height),
        "has_text": bool(getattr(shape, "has_text_frame", False)),
        "text": shape.text.strip() if getattr(shape, "has_text_frame", False) else "",
    }


def inspect_template(path: Path) -> dict[str, Any]:
    prs = Presentation(str(path))
    payload: dict[str, Any] = {
        "path": str(path),
        "slide_width_in": _emu(prs.slide_width),
        "slide_height_in": _emu(prs.slide_height),
        "slide_count": len(prs.slides),
        "layouts": [],
    }
    for index, layout in enumerate(prs.slide_layouts):
        shapes = [_shape_payload(shape) for shape in layout.shapes]
        payload["layouts"].append(
            {
                "index": index,
                "name": layout.name,
                "placeholder_count": sum(1 for shape in shapes if shape["is_placeholder"]),
                "text_shape_count": sum(1 for shape in shapes if shape["has_text"]),
                "shapes": shapes,
            }
        )
    return payload


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("template", type=Path)
    parser.add_argument("--output", type=Path)
    args = parser.parse_args()

    payload = inspect_template(args.template)
    text = json.dumps(payload, ensure_ascii=False, indent=2)
    if args.output:
        args.output.parent.mkdir(parents=True, exist_ok=True)
        args.output.write_text(text + "\n", encoding="utf-8")
    else:
        print(text)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

